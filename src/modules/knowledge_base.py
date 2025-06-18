import os
import json
import numpy as np
from pathlib import Path
from typing import List, Dict, Any, Optional

from src.modules.knowledge_base.storage.vector_store import VectorStore

class KnowledgeBase:
    """知识库管理模块，支持多种格式文档导入和语义检索"""
    
    def __init__(self, embedding_model: str = "bge-m3"):
        """
        初始化知识库
        
        Args:
            embedding_model: 嵌入模型名称，用于文本向量化
        """
        self.documents = []
        self.embeddings = []
        self.embedding_model = embedding_model
        self.vector_store = None
        self.metadata = {}
        
        # 支持的文件类型及对应的处理方法
        self.supported_formats = {
            ".pdf": self._process_pdf,
            ".txt": self._process_text,
            ".md": self._process_markdown,
            ".docx": self._process_docx,
            ".pptx": self._process_pptx,
            ".xlsx": self._process_excel,
            ".py": self._process_code,
            ".js": self._process_code,
            ".java": self._process_code,
            ".cpp": self._process_code,
            ".c": self._process_code,
            ".html": self._process_html,
        }
        
        print(f"知识库初始化完成，使用嵌入模型: {embedding_model}")
    
    def init_vector_store(self, kb_name: str, dimension: int = 768) -> bool:
        """
        初始化向量存储
        
        Args:
            kb_name: 知识库名称
            dimension: 向量维度
            
        Returns:
            bool: 是否初始化成功
        """
        try:
            # 设置向量存储路径
            db_path = Path("data") / "knowledge_bases" / kb_name / "vector_store.db"
            os.makedirs(db_path.parent, exist_ok=True)
            
            # 初始化向量存储
            self.vector_store = VectorStore(str(db_path), dimension=dimension)
            
            # 更新元数据
            self.metadata["kb_name"] = kb_name
            self.metadata["vector_dimension"] = dimension
            self.metadata["created_at"] = str(Path(db_path).stat().st_ctime) if Path(db_path).exists() else ""
            
            print(f"向量存储初始化成功: {db_path}")
            return True
        except Exception as e:
            print(f"初始化向量存储失败: {str(e)}")
            return False
    
    def add_document(self, file_path: str) -> bool:
        """
        添加文档到知识库
        
        Args:
            file_path: 文档路径
            
        Returns:
            bool: 是否添加成功
        """
        path = Path(file_path)
        if not path.exists():
            print(f"错误: 文件不存在 - {file_path}")
            return False
            
        ext = path.suffix.lower()
        if ext not in self.supported_formats:
            print(f"错误: 不支持的文件格式 - {ext}")
            return False
            
        # 调用对应格式的处理方法
        processor = self.supported_formats[ext]
        success = processor(file_path)
        
        if success:
            print(f"成功添加文档: {path.name}")
            return True
        else:
            print(f"添加文档失败: {path.name}")
            return False
    
    def add_web_content(self, url: str) -> bool:
        """
        添加网页内容到知识库
        
        Args:
            url: 网页URL
            
        Returns:
            bool: 是否添加成功
        """
        try:
            import requests
            from bs4 import BeautifulSoup

            print(f"正在获取网页内容: {url}")
            response = requests.get(url, timeout=10)
            response.raise_for_status()
            
            # 解析HTML内容
            soup = BeautifulSoup(response.text, 'html.parser')
            
            # 提取正文内容
            # 移除脚本、样式等无关内容
            for script in soup(["script", "style", "nav", "footer", "header"]):
                script.extract()
            
            # 获取文本内容
            text = soup.get_text(separator="\n", strip=True)
            
            # 添加到知识库
            success = self.add_text(text, {"source": url, "type": "web_content", "title": soup.title.text if soup.title else ""})
            return success
        except Exception as e:
            print(f"获取网页内容失败: {str(e)}")
            return False
    
    def add_text(self, text: str, metadata: Dict[str, Any] = None) -> bool:
        """
        直接添加文本内容到知识库
        
        Args:
            text: 文本内容
            metadata: 元数据
            
        Returns:
            bool: 是否添加成功
        """
        if not text.strip():
            return False
        
        if self.vector_store is None:
            print("警告: 向量存储未初始化，将只存储在内存中")
            
        # 对文本进行分块处理
        chunks = self._chunk_text(text, chunk_size=500, overlap=50)
        
        # 向量化处理
        try:
            # 生成文本嵌入向量
            embeddings = self._get_embeddings(chunks)
            
            # 如果向量存储已初始化，则添加到向量存储
            if self.vector_store is not None and embeddings:
                # 准备元数据
                metadatas = []
                for i, chunk in enumerate(chunks):
                    chunk_metadata = {
                        "chunk_id": f"{len(self.documents) + i}",
                        "chunk_index": i,
                        "total_chunks": len(chunks)
                    }
                    if metadata:
                        chunk_metadata.update(metadata)
                    metadatas.append(chunk_metadata)
                
                # 添加到向量存储
                self.vector_store.add(chunks, embeddings, metadatas)
                print(f"已添加到向量存储，{len(chunks)} 个块")
            
            # 同时存储在内存中
            for i, (chunk, embedding) in enumerate(zip(chunks, embeddings)):
                chunk_metadata = {
                    "chunk_id": f"{len(self.documents) + i}",
                    "chunk_index": i,
                    "total_chunks": len(chunks)
                }
                if metadata:
                    chunk_metadata.update(metadata)
                
                self.documents.append({"text": chunk, "metadata": chunk_metadata})
                self.embeddings.append(embedding)
            
            print(f"已添加文本内容，分为 {len(chunks)} 个块")
            return True
        except Exception as e:
            print(f"添加文本内容失败: {str(e)}")
            # 降级处理: 如果向量化失败，仍然存储原始文本
            for i, chunk in enumerate(chunks):
                chunk_metadata = {
                    "chunk_id": f"{len(self.documents) + i}",
                    "chunk_index": i,
                    "total_chunks": len(chunks),
                    "vectorized": False
                }
                if metadata:
                    chunk_metadata.update(metadata)
                
                self.documents.append({"text": chunk, "metadata": chunk_metadata})
            
            print(f"已添加文本内容(无向量)，分为 {len(chunks)} 个块")
            return True
    
    def search(self, query: str, top_k: int = 5) -> List[Dict[str, Any]]:
        """
        搜索知识库
        
        Args:
            query: 搜索查询
            top_k: 返回结果数量
            
        Returns:
            List[Dict]: 搜索结果列表
        """
        if self.vector_store is None and not self.documents:
            print("知识库为空，无法搜索")
            return []
            
        print(f"搜索: {query}")
        results = []
        
        try:
            # 生成查询向量
            query_embedding = self._get_embeddings([query])[0]
            
            # 优先使用向量存储进行搜索
            if self.vector_store is not None:
                vector_results = self.vector_store.search(query_embedding, k=top_k)
                for doc_id, score, metadata in vector_results:
                    results.append({
                        "text": metadata.get("text", ""),
                        "metadata": metadata,
                        "score": score
                    })
            
            # 如果向量存储结果不足，或者未使用向量存储，则使用内存中的文档
            if not results and self.embeddings:
                # 计算相似度
                similarities = []
                for i, doc_embedding in enumerate(self.embeddings):
                    similarity = self._calculate_similarity(query_embedding, doc_embedding)
                    similarities.append((i, similarity))
                
                # 排序并获取前K个结果
                similarities.sort(key=lambda x: x[1], reverse=True)
                for i, score in similarities[:top_k]:
                    results.append({
                        "text": self.documents[i]["text"],
                        "metadata": self.documents[i]["metadata"],
                        "score": float(score)
                    })
            
            # 如果没有向量，则返回最新添加的文档
            if not results and self.documents:
                print("警告: 无法进行向量搜索，返回最近添加的文档")
                for i in range(min(top_k, len(self.documents))):
                    idx = len(self.documents) - i - 1  # 从最新文档开始
                    results.append({
                        "text": self.documents[idx]["text"],
                        "metadata": self.documents[idx]["metadata"],
                        "score": 0.1  # 低置信度分数
                    })
                    
            return results
            
        except Exception as e:
            print(f"搜索失败: {str(e)}")
            # 降级处理: 返回最近添加的文档
            for i in range(min(top_k, len(self.documents))):
                idx = len(self.documents) - i - 1  # 从最新文档开始
                results.append({
                    "text": self.documents[idx]["text"],
                    "metadata": self.documents[idx]["metadata"],
                    "score": 0.1  # 低置信度分数
                })
            return results
    
    def load(self, directory: str) -> bool:
        """
        从目录加载知识库
        
        Args:
            directory: 知识库目录
            
        Returns:
            bool: 是否加载成功
        """
        kb_path = Path(directory)
        if not kb_path.exists() or not kb_path.is_dir():
            print(f"错误: 知识库目录不存在 - {directory}")
            return False
            
        # 加载元数据
        metadata_path = kb_path / "metadata.json"
        if metadata_path.exists():
            with open(metadata_path, 'r', encoding='utf-8') as f:
                self.metadata = json.load(f)
        
        # 加载向量索引        
        vector_db_path = kb_path / "vector_store.db"
        vector_index_path = kb_path / "vector_index.faiss"
        
        if vector_db_path.exists() and vector_index_path.exists():
            try:
                # 获取向量维度
                dimension = self.metadata.get("vector_dimension", 768)
                
                # 初始化向量存储并加载索引
                self.vector_store = VectorStore(str(vector_db_path), dimension=dimension)
                load_success = self.vector_store.load(str(vector_index_path))
                
                if load_success:
                    kb_name = self.metadata.get("kb_name", kb_path.name)
                    print(f"已加载知识库: {kb_name}, 包含 {self.vector_store.count()} 个文档")
                    return True
                else:
                    print("警告: 向量索引加载失败")
            except Exception as e:
                print(f"加载向量存储失败: {str(e)}")
        
        print(f"已加载知识库: {directory} (无向量索引)")
        return True
    
    def save(self, directory: str) -> bool:
        """
        保存知识库到目录
        
        Args:
            directory: 目标目录
            
        Returns:
            bool: 是否保存成功
        """
        kb_path = Path(directory)
        os.makedirs(kb_path, exist_ok=True)
        
        # 更新元数据
        self.metadata["document_count"] = len(self.documents)
        self.metadata["embedding_model"] = self.embedding_model
        
        # 保存元数据
        with open(kb_path / "metadata.json", 'w', encoding='utf-8') as f:
            json.dump(self.metadata, f, ensure_ascii=False, indent=2)
            
        # 保存向量索引
        if self.vector_store is not None:
            vector_index_path = kb_path / "vector_index.faiss"
            save_success = self.vector_store.save(str(vector_index_path))
            
            if save_success:
                print(f"已保存向量索引: {vector_index_path}")
            else:
                print("警告: 向量索引保存失败")
        
        print(f"已保存知识库: {directory}")
        return True
    
    def _get_embeddings(self, texts: List[str]) -> List[List[float]]:
        """
        获取文本嵌入向量
        
        Args:
            texts: 文本列表
            
        Returns:
            List[List[float]]: 嵌入向量列表
        """
        # 实际项目中，这里会调用嵌入模型API
        # 这里使用随机向量作为示例
        try:
            from src.services import model_manager
            
            # 尝试使用模型服务获取嵌入向量
            model = model_manager.get_embedding_model()
            if model:
                return model.get_embeddings(texts)
            
            print("警告: 无法获取嵌入模型，使用随机向量")
        except Exception as e:
            print(f"获取嵌入向量失败: {str(e)}")
        
        # 降级处理: 生成随机向量
        dim = 768  # 使用默认维度
        return [np.random.randn(dim).tolist() for _ in texts]
    
    def _calculate_similarity(self, vec1: List[float], vec2: List[float]) -> float:
        """计算两个向量的余弦相似度"""
        if not vec1 or not vec2:
            return 0.0
            
        vec1 = np.array(vec1)
        vec2 = np.array(vec2)
        
        norm1 = np.linalg.norm(vec1)
        norm2 = np.linalg.norm(vec2)
        
        if norm1 == 0 or norm2 == 0:
            return 0.0
            
        return np.dot(vec1, vec2) / (norm1 * norm2)
    
    # 文本分块处理
    def _chunk_text(self, text: str, chunk_size: int = 500, overlap: int = 50) -> List[str]:
        """
        将文本分割成块
        
        Args:
            text: 要分块的文本
            chunk_size: 每块的大致字符数
            overlap: 块之间的重叠字符数
            
        Returns:
            List[str]: 文本块列表
        """
        if not text:
            return []
            
        # 如果文本长度小于块大小，直接返回
        if len(text) <= chunk_size:
            return [text]
            
        chunks = []
        start = 0
        
        while start < len(text):
            # 确定块的结束位置
            end = min(start + chunk_size, len(text))
            
            # 尝试在句子边界处截断
            if end < len(text):
                # 查找最后一个句子结束符
                for sep in ['. ', '。', '! ', '？', '\n\n']:
                    last_sep = text[start:end].rfind(sep)
                    if last_sep != -1:
                        end = start + last_sep + len(sep)
                        break
                        
            # 添加当前块
            chunks.append(text[start:end])
            
            # 更新起始位置（考虑重叠）
            start = max(start + 1, end - overlap)
            
        return chunks
    
    # 各种文件格式的处理方法
    def _process_pdf(self, file_path: str) -> bool:
        """处理PDF文件"""
        try:
            import PyPDF2
            
            with open(file_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                text_content = ""
                
                # 提取每一页的文本
                for page_num in range(len(reader.pages)):
                    page = reader.pages[page_num]
                    text_content += page.extract_text() + "\n\n"
                    
                # 创建元数据
                metadata = {
                    "source": Path(file_path).name,
                    "file_path": str(file_path),
                    "type": "pdf",
                    "pages": len(reader.pages)
                }
                
                # 添加到知识库
                return self.add_text(text_content, metadata)
                
        except Exception as e:
            print(f"处理PDF文件失败: {str(e)}")
            return False
        
    def _process_text(self, file_path: str) -> bool:
        """处理文本文件"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                content = file.read()
                
            # 创建元数据
            metadata = {
                "source": Path(file_path).name,
                "file_path": str(file_path),
                "type": "text"
            }
            
            # 添加到知识库
            return self.add_text(content, metadata)
            
        except Exception as e:
            print(f"处理文本文件失败: {str(e)}")
            return False
        
    def _process_markdown(self, file_path: str) -> bool:
        """处理Markdown文件"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                content = file.read()
                
            # 创建元数据
            metadata = {
                "source": Path(file_path).name,
                "file_path": str(file_path),
                "type": "markdown"
            }
            
            # 尝试提取标题
            import re
            title_match = re.search(r'^#\s+(.+)$', content, re.MULTILINE)
            if title_match:
                metadata["title"] = title_match.group(1)
                
            # 添加到知识库
            return self.add_text(content, metadata)
            
        except Exception as e:
            print(f"处理Markdown文件失败: {str(e)}")
            return False
        
    def _process_docx(self, file_path: str) -> bool:
        """处理Word文档"""
        try:
            from docx import Document
            
            doc = Document(file_path)
            content = "\n".join([para.text for para in doc.paragraphs])
            
            # 创建元数据
            metadata = {
                "source": Path(file_path).name,
                "file_path": str(file_path),
                "type": "docx"
            }
            
            # 添加到知识库
            return self.add_text(content, metadata)
            
        except Exception as e:
            print(f"处理Word文档失败: {str(e)}")
            return False
        
    def _process_pptx(self, file_path: str) -> bool:
        """处理PowerPoint文件"""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            content = ""
            
            # 提取每个幻灯片的文本
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text + "\n\n"
            
            # 创建元数据
            metadata = {
                "source": Path(file_path).name,
                "file_path": str(file_path),
                "type": "pptx",
                "slides": len(prs.slides)
            }
            
            # 添加到知识库
            return self.add_text(content, metadata)
            
        except Exception as e:
            print(f"处理PowerPoint文件失败: {str(e)}")
            return False
        
    def _process_excel(self, file_path: str) -> bool:
        """处理Excel文件"""
        try:
            import pandas as pd
            
            # 读取所有工作表
            excel_file = pd.ExcelFile(file_path)
            content = ""
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(excel_file, sheet_name=sheet_name)
                content += f"工作表: {sheet_name}\n"
                content += df.to_string() + "\n\n"
            
            # 创建元数据
            metadata = {
                "source": Path(file_path).name,
                "file_path": str(file_path),
                "type": "excel",
                "sheets": len(excel_file.sheet_names)
            }
            
            # 添加到知识库
            return self.add_text(content, metadata)
            
        except Exception as e:
            print(f"处理Excel文件失败: {str(e)}")
            return False
        
    def _process_code(self, file_path: str) -> bool:
        """处理代码文件"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                content = file.read()
                
            # 确定编程语言
            ext = Path(file_path).suffix.lower()
            language_map = {
                ".py": "python",
                ".js": "javascript",
                ".ts": "typescript",
                ".java": "java",
                ".cpp": "cpp",
                ".c": "c",
                ".cs": "csharp",
                ".go": "go",
                ".rb": "ruby",
                ".php": "php",
                ".rs": "rust",
                ".swift": "swift",
                ".kt": "kotlin"
            }
            
            language = language_map.get(ext, "unknown")
            
            # 创建元数据
            metadata = {
                "source": Path(file_path).name,
                "file_path": str(file_path),
                "type": "code",
                "code_lang": language,
                "lines": content.count('\n') + 1
            }
            
            # 添加到知识库
            return self.add_text(content, metadata)
            
        except Exception as e:
            print(f"处理代码文件失败: {str(e)}")
            return False
        
    def _process_html(self, file_path: str) -> bool:
        """处理HTML文件"""
        try:
            from bs4 import BeautifulSoup
            
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                content = file.read()
            
            # 解析HTML
            soup = BeautifulSoup(content, 'html.parser')
            
            # 提取标题
            title = soup.title.string if soup.title else ""
            
            # 移除脚本、样式等内容
            for script in soup(["script", "style"]):
                script.extract()
                
            # 获取文本内容
            text_content = soup.get_text(separator="\n", strip=True)
            
            # 创建元数据
            metadata = {
                "source": Path(file_path).name,
                "file_path": str(file_path),
                "type": "html",
                "title": title
            }
            
            # 添加到知识库
            return self.add_text(text_content, metadata)
            
        except Exception as e:
            print(f"处理HTML文件失败: {str(e)}")
            return False 